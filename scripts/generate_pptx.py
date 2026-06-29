import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # Colors
    c_navy = RGBColor(12, 35, 64)
    c_blue = RGBColor(30, 144, 255)
    c_dark = RGBColor(50, 50, 50)
    c_gray = RGBColor(128, 128, 128)
    
    # Helper to style text frames
    def set_font(run, size=18, bold=False, italic=False, color=c_dark, font_name="Calibri"):
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color

    # Helper to add standard slide with Title
    def add_slide(title_text):
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)
        
        # Add Title Box
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Calibri"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = c_navy
        
        return slide

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Background Accent shape
    accent = slide1.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(4.5), Inches(7.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = c_navy
    accent.line.color.rgb = c_navy
    
    title_box = slide1.shapes.add_textbox(Inches(5.0), Inches(2.0), Inches(7.8), Inches(4.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "Reinforcement Learning for Underwater Acoustic Vessel Detection and Tracking"
    set_font(p.runs[0], size=36, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "\nCourse: Introduction to Reinforcement Learning"
    set_font(p2.runs[0], size=20, color=c_blue)
    
    p3 = tf.add_paragraph()
    p3.text = "Author: Roy Michael\nDate: June 2026"
    set_font(p3.runs[0], size=18, color=c_gray)
    
    notes1 = slide1.notes_slide.notes_text_frame
    notes1.text = (
        "Welcome everyone to this presentation on Reinforcement Learning for Underwater Acoustic "
        "Vessel Detection and Tracking.\n\n"
        "In this talk, we will discuss how to formulate passive sonar tracking as a Markov Decision "
        "Process, and evaluate Double Q-Learning, Linear Function Approximation (Tile Coding), "
        "and Actor-Critic policies on real-world ocean datasets collected in Haifa and Croatia."
    )

    # Slide 2: Raw Waveform vs. Spectrogram
    slide2 = add_slide("Acoustic Signals: From Time to Frequency")
    
    # Left Image (Raw Waveform)
    if os.path.exists("report/images/raw_acoustic_signal.png"):
        slide2.shapes.add_picture("report/images/raw_acoustic_signal.png", Inches(0.5), Inches(1.6), width=Inches(6.0))
        
    # Right Image (Spectrogram)
    if os.path.exists("report/images/LOFAR_Joint_Signal.png"):
        slide2.shapes.add_picture("report/images/LOFAR_Joint_Signal.png", Inches(6.8), Inches(1.6), width=Inches(6.0))
        
    tb = slide2.shapes.add_textbox(Inches(0.5), Inches(5.6), Inches(12.33), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.text = "• Raw Time-Domain: Acoustic pressure is heavily masked by ocean clutter (transients, waves), hiding targets.\n" \
              "• LOFAR Spectrogram (STFT): Converts signal to time-frequency domain, exposing vessel engine signatures as narrowband tonals."
    set_font(p1.runs[0], size=16, color=c_dark)

    notes2 = slide2.notes_slide.notes_text_frame
    notes2.text = (
        "Let's look at the transition from raw sound waves to the spectrogram.\n\n"
        "On the left, we see the raw time-domain pressure waveform recorded by our hydrophone. "
        "It's extremely noisy, with snapping shrimp clicks and wave noise completely masking the vessel.\n\n"
        "To make tracking possible, we compute the Short-Time Fourier Transform (STFT) to produce the LOFAR spectrogram "
        "shown on the right. This converts the signal into the frequency domain, where vessel engine tonals manifest as "
        "continuous horizontal lines."
    )

    # Slide 3: Physical Challenges in Passive Sonar Tracking
    slide3 = add_slide("Physical Challenges in Sonar Tracking")
    
    # Left Image (Annotated Spectrogram)
    if os.path.exists("report/images/lofar_annotated.png"):
        slide3.shapes.add_picture("report/images/lofar_annotated.png", Inches(0.5), Inches(1.5), width=Inches(5.8))
        
    # Right Text Column
    tb = slide3.shapes.add_textbox(Inches(6.6), Inches(1.5), Inches(6.2), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Tracking Obstacles in Real Ocean Channels:"
    set_font(p1.runs[0], size=20, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• Doppler Frequency Drift:"
    set_font(p2.runs[0], size=18, bold=True, color=c_navy)
    p2.space_before = Pt(10)
    p3 = tf.add_paragraph()
    p3.text = "  - Vessel acceleration shifts observed tones, causing lines to 'slide' across frequency bins."
    set_font(p3.runs[0], size=15, color=c_dark)
    
    p4 = tf.add_paragraph()
    p4.text = "• Multipath Reflections & Echoes:"
    set_font(p4.runs[0], size=18, bold=True, color=c_navy)
    p4.space_before = Pt(10)
    p5 = tf.add_paragraph()
    p5.text = "  - Surface/seabed reflections create phantom tracks and parallel harmonic overtones (shown in red/orange)."
    set_font(p5.runs[0], size=15, color=c_dark)
    
    p6 = tf.add_paragraph()
    p6.text = "• Acoustic Fading (Doppler Nulls):"
    set_font(p6.runs[0], size=18, bold=True, color=c_navy)
    p6.space_before = Pt(10)
    p7 = tf.add_paragraph()
    p7.text = "  - Destructive interference fades signals temporarily, requiring memory to prevent dropping active tracks."
    set_font(p7.runs[0], size=15, color=c_dark)

    notes3 = slide3.notes_slide.notes_text_frame
    notes3.text = (
        "Passive sonar tracking is highly complex due to physical underwater phenomena.\n\n"
        "First, Doppler shift causes frequencies to slide as vessels accelerate. "
        "Second, multipath reflections create duplicate lines and harmonics, which are highlighted on our spectrogram graph. "
        "Third, acoustic fading can make a vessel tone completely disappear for several seconds. "
        "Our RL agent must learn to maintain track state through these fades and Doppler drifts."
    )

    # Slide 4: Non-negative Matrix Factorization (NMF)
    slide4 = add_slide("Non-negative Matrix Factorization (NMF)")
    
    # Left Text Column
    tb = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Isolating Targets from Ambient Noise:"
    set_font(p1.runs[0], size=20, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• Mathematical Decompostion (V ~ H * W):"
    set_font(p2.runs[0], size=18, bold=True, color=c_navy)
    p2.space_before = Pt(15)
    p3 = tf.add_paragraph()
    p3.text = "  - Factorizes the spectrogram into static spectral profiles (H) and temporal activation weights (W)."
    set_font(p3.runs[0], size=16, color=c_dark)
    
    p4 = tf.add_paragraph()
    p4.text = "• Pre-Processing Filter for RL:"
    set_font(p4.runs[0], size=18, bold=True, color=c_navy)
    p4.space_before = Pt(15)
    p5 = tf.add_paragraph()
    p5.text = "  - Extracts clean, stable narrowband engine components while filtering out random broadband clutter (waves, cavitation)."
    set_font(p5.runs[0], size=16, color=c_dark)
    
    # Right Image Column
    if os.path.exists("report/images/NMF_Components_Joint_Signal.png"):
        slide4.shapes.add_picture("report/images/NMF_Components_Joint_Signal.png", Inches(6.8), Inches(1.8), width=Inches(6.0))
        
    notes4 = slide4.notes_slide.notes_text_frame
    notes4.text = (
        "Non-negative Matrix Factorization acts as our unsupervised pre-processing filter.\n\n"
        "By decomposing the spectrogram V into component dictionary profiles H and activations W, we successfully "
        "isolate the steady engine tonals from broadband cavitation and wave noise. "
        "This allows the DSP orchestrator to extract clean candidate peaks for the reinforcement learning tracker."
    )

    # Slide 6: Project Overview & Motivation
    slide6 = add_slide("Project Overview & Motivation")
    tb = slide6.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "• The Challenge of Passive Sonar Tracking:"
    set_font(p1.runs[0], size=20, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "  - Marine environments are filled with non-stationary background noise (wave clutter, biological clicks)."
    set_font(p2.runs[0], size=18, color=c_dark)
    
    p3 = tf.add_paragraph()
    p3.text = "  - Vessel acceleration and maneuvers induce large frequency shifts (Doppler drifts) and temporary fading."
    set_font(p3.runs[0], size=18, color=c_dark)
    
    p4 = tf.add_paragraph()
    p4.text = "• Limitations of Heuristic Trackers:"
    set_font(p4.runs[0], size=20, bold=True, color=c_navy)
    p4.space_before = Pt(20)
    
    p5 = tf.add_paragraph()
    p5.text = "  - Traditional rule-based trackers use rigid gates. Narrow gates lose tracks during acceleration; wide gates associate with noise."
    set_font(p5.runs[0], size=18, color=c_dark)
    
    p6 = tf.add_paragraph()
    p6.text = "• RL Formulation Goal:"
    set_font(p6.runs[0], size=20, bold=True, color=c_navy)
    p6.space_before = Pt(20)
    
    p7 = tf.add_paragraph()
    p7.text = "  - Automate tracking decisions as an MDP, balancing track continuity with noise rejection to maximize long-term rewards."
    set_font(p7.runs[0], size=18, color=c_dark)

    notes6 = slide6.notes_slide.notes_text_frame
    notes6.text = (
        "Here we introduce the motivation for our work.\n\n"
        "Passive sonar processing is inherently difficult because the underwater channel is highly non-stationary. "
        "We have biological noises, snapping shrimp, wave action, and wind. Traditional systems use heuristically tuned "
        "gating thresholds. If the vessel speeds up, its engine tone slides in frequency due to the Doppler shift. A narrow "
        "threshold will immediately drop the track, while a wide threshold will associate with nearby clutter.\n\n"
        "We want to reformulate this as a sequential decision-making process where the agent learns the optimal policy "
        "under uncertainty."
    )

    # Slide 7: Experimental Setup
    slide7 = add_slide("Field Experiments & Mooring Setup")
    
    # Left text column
    tb = slide7.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Croatia Ocean Sonics Dataset (Silba Island):"
    set_font(p1.runs[0], size=18, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• Radiated noise from Suex VRX Diver Propulsion Vehicles (DPVs) and surface support boats.\n• Real-world interference including passing ferries."
    set_font(p2.runs[0], size=16, color=c_dark)
    
    p3 = tf.add_paragraph()
    p3.text = "Haifa Dataset:"
    set_font(p3.runs[0], size=18, bold=True, color=c_navy)
    p3.space_before = Pt(15)
    
    p4 = tf.add_paragraph()
    p4.text = "• Seacraft GO! DVP executing circular and linear maneuvers."
    set_font(p4.runs[0], size=16, color=c_dark)
    
    p5 = tf.add_paragraph()
    p5.text = "Deployment Equipment:"
    set_font(p5.runs[0], size=18, bold=True, color=c_navy)
    p5.space_before = Pt(15)
    
    p6 = tf.add_paragraph()
    p6.text = "• icListen HF ALTA hydrophone unit (128 kHz sampling rate).\n• Moored on seabed frames at 10-30m depths to eliminate surface waves and buoy movements."
    set_font(p6.runs[0], size=16, color=c_dark)
    
    # Right Image column
    if os.path.exists("report/images/silba-1k.png"):
        slide7.shapes.add_picture("report/images/silba-1k.png", Inches(7.0), Inches(1.8), width=Inches(5.5))
        
    notes7 = slide7.notes_slide.notes_text_frame
    notes7.text = (
        "To gather realistic acoustic signatures, field trials were performed in Haifa, Israel and Silba Island, Croatia.\n\n"
        "The targets of interest are Diver Propulsion Vehicles, also known as underwater scooters, which emit quiet, "
        "low-frequency narrowband motor tonals from their pulse-width modulated electric drives. Support boats and ferries "
        "are also recorded, representing typical harbor noise clutter.\n\n"
        "Crucially, the hydrophone is placed on a seabed-anchored frame. This isolates the hydrophone from the movement of the "
        "surface waves and cable strumming, which would otherwise introduce massive low-frequency noise."
    )

    # Slide 13: RL-DSP System Architecture & Packages
    slide13 = add_slide("RL-DSP System Architecture & Packages")
    
    # Left text column
    tb = slide13.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Separation of Concerns:"
    set_font(p1.runs[0], size=18, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• core.environment Package:"
    set_font(p2.runs[0], size=16, bold=True, color=c_navy)
    p2.space_before = Pt(8)
    p3 = tf.add_paragraph()
    p3.text = "  - AcousticDataStreamer streams spectral frames.\n  - TrackingMDPEnv wraps streamer, computes state, resolves rewards."
    set_font(p3.runs[0], size=14, color=c_dark)
    
    p4 = tf.add_paragraph()
    p4.text = "• core.agent Package:"
    set_font(p4.runs[0], size=16, bold=True, color=c_navy)
    p4.space_before = Pt(12)
    p5 = tf.add_paragraph()
    p5.text = "  - DSPOrchestrator manages NMF updates, extracts peaks, routes peaks to dynamic trackers.\n  - VesselTrackProcessor tracks independent vessels."
    set_font(p5.runs[0], size=14, color=c_dark)
    
    # Right Image column
    if os.path.exists("report/images/high_level_architecture.png"):
        slide13.shapes.add_picture("report/images/high_level_architecture.png", Inches(6.5), Inches(1.8), width=Inches(6.3))

    notes13 = slide13.notes_slide.notes_text_frame
    notes13.text = (
        "Here we look at the system architecture.\n\n"
        "We divide the code strictly into two packages: core.environment and core.agent.\n\n"
        "In the environment package, we run the low-level signal processing like the STFT frame buffer and the standard Gym-like "
        "MDP wrapper. The agent package houses the orchestrator, which runs the NMF model on the STFT buffer, extracts candidate "
        "peak frequencies, and manages dynamically spawned vessel track processors.\n\n"
        "This boundary isolates the deterministic physical processing (NMF and clustering) from the stochastically updated "
        "reinforcement learning policy."
    )

    # Slide 13b: System Dynamics & Component Interaction
    slide_dyn = add_slide("System Dynamics & Component Interaction")
    
    # Left text column
    tb = slide_dyn.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Detailed Processing Dynamics:"
    set_font(p1.runs[0], size=18, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• 1. Acoustic Feed:\n" \
             "  - Spectrogram peaks are extracted by the NMF engine."
    set_font(p2.runs[0], size=14, color=c_dark)
    p2.space_before = Pt(8)
    
    p3 = tf.add_paragraph()
    p3.text = "• 2. State Query:\n" \
             "  - The Orchestrator constructs state S (nearest distance, amplitude, tonality, age) and queries the Agent."
    set_font(p3.runs[0], size=14, color=c_dark)
    p3.space_before = Pt(8)
    
    p4 = tf.add_paragraph()
    p4.text = "• 3. Policy Action:\n" \
             "  - The RL Agent returns action A (Spawn, Associate, or Reject)."
    set_font(p4.runs[0], size=14, color=c_dark)
    p4.space_before = Pt(8)
    
    p5 = tf.add_paragraph()
    p5.text = "• 4. Execution:\n" \
             "  - The Orchestrator routes the peak to update concurrent track processors."
    set_font(p5.runs[0], size=14, color=c_dark)
    p5.space_before = Pt(8)

    # Right Image column
    if os.path.exists("report/images/architecture_diagram.png"):
        slide_dyn.shapes.add_picture("report/images/architecture_diagram.png", Inches(6.6), Inches(1.8), width=Inches(6.2))

    notes_dyn = slide_dyn.notes_slide.notes_text_frame
    notes_dyn.text = (
        "This slide visualizes the system dynamics.\n\n"
        "Spectrogram peaks extracted by the NMF engine are handled by the DSPOrchestrator. "
        "The Orchestrator constructs state S and queries the RL Agent, which selects an action. "
        "Then the Orchestrator executes this action, updating the concurrent VesselTrackProcessor instances."
    )

    # Slide 8: MDP Formulation & State Space
    slide8 = add_slide("Problem Formulation: MDP State Space")
    tb = slide8.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Markov Decision Process (MDP) Tuple: <S, A, P, R, gamma>"
    set_font(p1.runs[0], size=22, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "Continuous State Space (S):"
    set_font(p2.runs[0], size=20, bold=True, color=c_navy)
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "  s_continuous = (d_Hz, A, T, age)"
    set_font(p3.runs[0], size=18, bold=True, color=c_blue)
    
    p4 = tf.add_paragraph()
    p4.text = "  • d_Hz : Spectral distance from the candidate peak to the closest active vessel track.\n" \
             "  • A : Relative amplitude (activation weight) of the spectral peak.\n" \
             "  • T : Tonality score (spectral flatness-based). Distinguishes tonals from noise.\n" \
             "  • age : Continuous track age of the closest active track."
    set_font(p4.runs[0], size=16, color=c_dark)
    
    p5 = tf.add_paragraph()
    p5.text = "Discretized State Bins (S_discrete) for Tabular Policies:"
    set_font(p5.runs[0], size=20, bold=True, color=c_navy)
    p5.space_before = Pt(10)
    
    p6 = tf.add_paragraph()
    p6.text = "  • 6 Distance Bins | 5 Amplitude Bins | 5 Tonality Bins | 4 Track Age Bins\n" \
             "  • Total State Space Size: 6 * 5 * 5 * 4 = 600 discrete states."
    set_font(p6.runs[0], size=16, color=c_dark)

    notes8 = slide8.notes_slide.notes_text_frame
    notes8.text = (
        "We formalize the tracking problem as an MDP.\n\n"
        "For each incoming acoustic peak, the state tuple s consists of distance, amplitude, tonality, and track age. "
        "The distance measures closeness to existing tracks. The amplitude and tonality represent signal quality. "
        "The track age helps separate temporary transients from established tracks.\n\n"
        "For tabular policies, we discretize these continuous dimensions into 6x5x5x4 bins, yielding 600 discrete states. "
        "This state representation is highly compact, allowing Q-Learning to converge very quickly while retaining the physical properties of the tracks."
    )

    # Slide 10: Markov Decision Process: Action Space
    slide10 = add_slide("Markov Decision Process: Action Space")
    tb = slide10.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Action Options (a) for Local Routing Decisions:"
    set_font(p1.runs[0], size=22, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• REJECT (a = 0): Discard the peak observation."
    set_font(p2.runs[0], size=18, bold=True, color=c_navy)
    p2.space_before = Pt(15)
    p3 = tf.add_paragraph()
    p3.text = "  - The orchestrator takes no action, letting inactive tracking states age."
    set_font(p3.runs[0], size=16, color=c_dark)
    
    p4 = tf.add_paragraph()
    p4.text = "• ASSOCIATE (a = 1): Route peak to closest active tracker."
    set_font(p4.runs[0], size=18, bold=True, color=c_navy)
    p4.space_before = Pt(15)
    p5 = tf.add_paragraph()
    p5.text = "  - Routes peak coordinates to update the coordinates of the closest active tracker."
    set_font(p5.runs[0], size=16, color=c_dark)
    
    p6 = tf.add_paragraph()
    p6.text = "• SPAWN (a = 2): Start a new vessel tracker."
    set_font(p6.runs[0], size=18, bold=True, color=c_navy)
    p6.space_before = Pt(15)
    p7 = tf.add_paragraph()
    p7.text = "  - Instantiates a new concurrent vessel tracker instance at the peak's frequency."
    set_font(p7.runs[0], size=16, color=c_dark)

    notes10 = slide10.notes_slide.notes_text_frame
    notes10.text = (
        "Here we show the core action options of the MDP.\n\n"
        "At each frame, the agent receives a peak and chooses one of three actions: Reject, Associate, or Spawn. "
        "Reject ignores the peak, Associate routes it to the closest active track, and Spawn initializes a new track."
    )

    # Slide 10b: Markov Decision Process: Decoupled Reward Rules
    slide_rew = add_slide("Markov Decision Process: Decoupled Reward Rules")
    tb = slide_rew.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Critic's Dense Feedback Calculations:"
    set_font(p1.runs[0], size=22, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• REJECT (a = 0): Correct Reject (+2.0) | False Negative Miss (-10.0)"
    set_font(p2.runs[0], size=16, color=c_dark)
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "• ASSOCIATE (a = 1): Good Assoc (+10.0, <=30Hz) | Speed Stage Change (+5.0, 30-65Hz) | Bad Mismatch (-15.0, >65Hz) | Invalid Assoc (-20.0)"
    set_font(p3.runs[0], size=16, color=c_dark)
    p3.space_before = Pt(10)
    
    p4 = tf.add_paragraph()
    p4.text = "• SPAWN (a = 2): Correct Spawn (+10.0/High, +5.0/Med) | Duplicate Spawn (-10.0, <=35Hz)"
    set_font(p4.runs[0], size=16, color=c_dark)
    p4.space_before = Pt(10)

    notes_rew = slide_rew.notes_slide.notes_text_frame
    notes_rew.text = (
        "This slide defines the reward rules computed by the environment critic. "
        "Correct actions yield positive rewards (like Good Association at +10.0 or Correct Reject at +2.0), "
        "while incorrect actions are penalized (such as duplicate spawns at -10.0 or mismatches at -15.0)."
    )

    # Slide 11: Rationale for Reward Magnitudes
    slide11 = add_slide("Reward Design: Preventing Degenerate Policies")
    tb = slide11.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Reward engineering was utilized to avoid critical exploit policies:"
    set_font(p1.runs[0], size=22, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "1. Preventing the 'Lazy' Agent Policy:"
    set_font(p2.runs[0], size=18, bold=True, color=c_navy)
    p2.space_before = Pt(15)
    p3 = tf.add_paragraph()
    p3.text = "   - If the reward for a Correct Reject was too high, the agent would learn a degenerate policy of rejecting every peak to safely accumulate points, avoiding any risk of association penalties.\n" \
             "   - Solution: Correct Reject reward is a small baseline (+2.0) while Good Association is five times higher (+10.0)."
    set_font(p3.runs[0], size=16, color=c_dark)
    
    p4 = tf.add_paragraph()
    p4.text = "2. Preventing 'Spawn Farming' Exploits:"
    set_font(p4.runs[0], size=18, bold=True, color=c_navy)
    p4.space_before = Pt(15)
    p5 = tf.add_paragraph()
    p5.text = "   - If spawning a tracker yielded equal or higher reward than maintaining it, the agent would constantly drop active tracks and spawn new ones to harvest points.\n" \
             "   - Solution: Moderate Spawn reward is capped at +5.0, keeping it lower than Good Association (+10.0)."
    set_font(p5.runs[0], size=16, color=c_dark)

    notes11 = slide11.notes_slide.notes_text_frame
    notes11.text = (
        "Designing the reward system in reinforcement learning is often a process of trial and error because agents are "
        "extremely good at finding loopholes or degenerate policies.\n\n"
        "For example, when we first set the rewards, the agent learned a 'lazy policy.' Because associating a peak carries the "
        "risk of a -15.0 mismatch penalty, and rejecting noise gives a positive reward, the agent realized it could maximize its "
        "score by simply rejecting every peak it saw, never starting any tracks. We solved this by scaling the association "
        "reward to be five times larger than the reject reward.\n\n"
        "Similarly, we had to moderate the spawn reward to prevent 'spawn farming' where the agent would spawn a tracker and "
        "immediately drop it to spawn another."
    )

    # Slide 14: Evaluated Algorithms
    slide14 = add_slide("Evaluated Reinforcement Learning Algorithms")
    tb = slide14.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Paradigm Selection & Empirical Performance:"
    set_font(p1.runs[0], size=22, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "1. Tabular Double Q-Learning:"
    set_font(p2.runs[0], size=18, bold=True, color=c_navy)
    p2.space_before = Pt(10)
    p3 = tf.add_paragraph()
    p3.text = "   - Selection: Simple, computationally light baseline to validate discrete state space; resolves maximization bias using two independent Q-tables.\n   - Behavior: Converges quickly but suffers from track-stitching gaps due to discretization limits."
    set_font(p3.runs[0], size=15, color=c_dark)
    
    p4 = tf.add_paragraph()
    p4.text = "2. Linear Function Approximation (Tile Coding):"
    set_font(p4.runs[0], size=18, bold=True, color=c_navy)
    p4.space_before = Pt(10)
    p5 = tf.add_paragraph()
    p5.text = "   - Selection: Handle continuous coordinates directly by mapping states to 5,184-dimensional sparse binary space via 4 overlapping tilings (each with 6^4 tiles).\n   - Behavior: Performed the best, achieving the highest cumulative rewards because it represents coordinate metrics continuously."
    set_font(p5.runs[0], size=15, color=c_dark)
    
    p6 = tf.add_paragraph()
    p6.text = "3. Actor-Critic (Stochastic Policy):"
    set_font(p6.runs[0], size=18, bold=True, color=c_navy)
    p6.space_before = Pt(10)
    p7 = tf.add_paragraph()
    p7.text = "   - Selection: Evaluates policy gradients on continuous features to support probabilistic associations under fading conditions.\n   - Behavior: Suffered from early preference saturation due to large reward bounds. Behaved stably only after introducing reward scaling, gradient clipping, and preference centering."
    set_font(p7.runs[0], size=15, color=c_dark)

    notes14 = slide14.notes_slide.notes_text_frame
    notes14.text = (
        "We benchmark three distinct algorithm classes representing tabular, continuous, and stochastic policy paradigms.\n\n"
        "Double Q-Learning serves as our robust tabular baseline, specifically selected to combat maximization bias. "
        "Linear Function Approximation uses multi-tiled coding to resolve continuous state coordinates, preventing boundary "
        "discretization issues. Actor-Critic models policy preferences directly via softmax, giving us a probabilistic tracker "
        "capable of maintaining tentative tracks during signal fading."
    )

    # Slide 15: Experimental Results - Croatia 2407_1 Metrics
    slide15 = add_slide("Experimental Results: Quantitative Metrics")
    
    # Left text / table column
    tb = slide15.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Performance Metrics (Croatia 2407_1 Dataset):"
    set_font(p1.runs[0], size=18, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• Linear FA achieves the highest cumulative reward (81,685) and good associations (8,815) due to continuous coordinate resolution.\n" \
             "• Double Q-Learning is extremely robust, achieving high track consistency with only 1.0% bad associations (reward: 78,595).\n" \
             "• Actor-Critic successfully resolves vessel tracks, achieving a cumulative reward of 74,531."
    set_font(p2.runs[0], size=14, color=c_dark)
    p2.space_before = Pt(15)
    
    # Add a simple text-based table representing the metrics
    p3 = tf.add_paragraph()
    p3.text = "Agent | Reward | Good Assoc | Bad Assoc % | Vessels\n" \
              "---------------------------------------------------\n" \
              "Double Q | 78,595 | 8,648 | 1.0% | 3\n" \
              "Linear FA | 81,685 | 8,815 | 1.3% | 3\n" \
              "Actor-Critic | 74,531 | 8,449 | 1.7% | 2"
    set_font(p3.runs[0], size=12, bold=True, font_name="Consolas", color=c_navy)
    p3.space_before = Pt(15)
    
    # Right Image column
    if os.path.exists("report/images/rl_comparison_croatia_2407_1.png"):
        slide15.shapes.add_picture("report/images/rl_comparison_croatia_2407_1.png", Inches(7.2), Inches(1.8), width=Inches(5.6))

    notes15 = slide15.notes_slide.notes_text_frame
    notes15.text = (
        "Here are the quantitative results. The table shows that Linear Function Approximation performs the best, "
        "accumulating the highest cumulative reward and most good associations. This is because continuous tile coding "
        "allows the agent to track fine-grained frequency slides without getting stuck on bin boundaries.\n\n"
        "Double Q-Learning, however, is a close second and remains highly robust, with only 2.1% bad associations. "
        "Actor-Critic also performs well, but has slightly more duplicate spawns due to its stochastic policy exploring "
        "actions in borderline states."
    )

    # Slide 16: Training Convergence Profile
    slide16 = add_slide("Training Convergence Profile")
    
    # Left Image column
    if os.path.exists("report/images/convergence_combined.png"):
        slide16.shapes.add_picture("report/images/convergence_combined.png", Inches(0.5), Inches(1.8), width=Inches(6.0))
        
    # Right Image column
    if os.path.exists("report/images/convergence_individual.png"):
        slide16.shapes.add_picture("report/images/convergence_individual.png", Inches(6.8), Inches(1.8), width=Inches(6.0))

    notes16 = slide16.notes_slide.notes_text_frame
    notes16.text = (
        "Let's look at the convergence graphs over 500 episodes.\n\n"
        "In early versions, a bug kept epsilon at 1.0, making training returns flat and negative. "
        "After correcting the bug to apply the decayed epsilon per episode, we see classic convergence: "
        "Double Q-Learning and Linear FA start at around -3,900 when exploration is 100% and climb steadily "
        "as epsilon decays down to 0.05, plateauing above +11,500.\n\n"
        "Actor-Critic begins and remains positive because its softmax policy stochastically explores actions, "
        "avoiding the uniform random penalties that trigger negative returns in early Q-learning episodes."
    )

    # Slide 17: Timeline and Normalization
    slide17 = add_slide("Absolute timelines & Policy Normalization")
    
    # Left text column
    tb = slide17.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Absolute Time timelines (HH:MM:SS):"
    set_font(p1.runs[0], size=18, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• Mapped directly to timestamps parsed from raw WAV filenames.\n" \
             "• Double Q-Learning resolves the main ferry transit and support boat tracks without identity swaps."
    set_font(p2.runs[0], size=14, color=c_dark)
    
    p3 = tf.add_paragraph()
    p3.text = "Actor-Critic Policy Saturation Resolve:"
    set_font(p3.runs[0], size=18, bold=True, color=c_navy)
    p3.space_before = Pt(15)
    
    p4 = tf.add_paragraph()
    p4.text = "Early AC flatlined due to massive TD-errors saturating the softmax function. We resolved this using:\n" \
             "• Reward Scaling: Scaling rewards down by 20.0.\n" \
             "• Gradient Clipping: Clipping TD-error within [-1.0, 1.0].\n" \
             "• Preference Centering: Subtracting mean preferences at each step to prevent weights from drifting to infinity."
    set_font(p4.runs[0], size=14, color=c_dark)
    
    # Right Image column
    if os.path.exists("report/images/croatia_2407_1_double_q_learning_timeline.png"):
        slide17.shapes.add_picture("report/images/croatia_2407_1_double_q_learning_timeline.png", Inches(6.8), Inches(1.8), width=Inches(6.0))

    notes17 = slide17.notes_slide.notes_text_frame
    notes17.text = (
        "Here we discuss two critical topics: absolute timelines and Actor-Critic normalization.\n\n"
        "First, we integrated absolute time tracking. By parsing the timestamps in the raw file names, we map our sonar "
        "detections directly to real-world time in HH:MM:SS. The chart on the right shows that Double Q-Learning reconstructs the "
        "boats' timelines perfectly, matching passing ferry events recorded in field logs.\n\n"
        "Second, we resolved Actor-Critic flatlining. Early on, unscaled rewards led to massive TD-errors, causing the softmax "
        "preferences to saturate (one action probability goes to 1.0). We solved this by scaling the rewards down by 20.0, "
        "clipping TD-errors, and centering actor preferences at each step."
    )

    # Slide 18: Conclusion
    slide18 = add_slide("Conclusions & Future Outlook")
    
    tb = slide18.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Key Contributions of this Project:"
    set_font(p1.runs[0], size=22, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "• Formulation of passive acoustic tracking as an autonomous, self-supervised MDP.\n" \
             "• Successful elimination of human-labeled training data requirements by utilizing physical consistency constraints.\n" \
             "• Benchmarking of tabular, continuous, and stochastic policy models under a unified codebase.\n" \
             "• Integration of absolute real-world timestamps for time-synchronized trajectory reconstruction."
    set_font(p2.runs[0], size=16, color=c_dark)
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Future Research Directions:"
    set_font(p3.runs[0], size=20, bold=True, color=c_navy)
    p3.space_before = Pt(20)
    
    p4 = tf.add_paragraph()
    p4.text = "• Reduce tile coding crosstalk (feature crosstalk interference) to stabilize continuous function approximation return profiles.\n" \
             "• Evaluate policy robustness under active acoustic jamming or extreme environmental surface clutter."
    set_font(p4.runs[0], size=16, color=c_dark)

    notes18 = slide18.notes_slide.notes_text_frame
    notes18.text = (
        "In conclusion, we have designed and validated a robust reinforcement learning framework for passive sonar vessel tracking.\n\n"
        "Our method eliminates the need for human labels, converting physical acoustic rules and temporal consistency into a "
        "self-supervised feedback loop.\n\n"
        "In the future, we plan to refine the continuous function approximation models to minimize feature crosstalk and evaluate "
        "robustness under active jamming or extreme marine storm noise."
    )

    # Slide 19: Demonstration Slide
    slide19 = add_slide("Interactive Playback & Demonstration")
    tb = slide19.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "How to run the live tracking animation:"
    set_font(p1.runs[0], size=22, bold=True, color=c_navy)
    
    p2 = tf.add_paragraph()
    p2.text = "1. Execute the orchestration evaluation script to stream frame updates:\n" \
              "   python run_orchestrator.py --train-dataset croatia_2407_1 --episodes 150 --eval-datasets croatia_2407_1\n" \
              "   (This automatically outputs chronological tracking JPEG frames to output/frames/)"
    set_font(p2.runs[0], size=16, color=c_dark)
    p2.space_before = Pt(10)
    
    p3 = tf.add_paragraph()
    p3.text = "2. Open viewer.html in a web browser to view playback:\n" \
              "   - The viewer runs a high-speed JavaScript slideshow (100ms interval) of the generated frames.\n" \
              "   - Shows real-time dynamic track spawning (a=2), sequential line associations (a=1), and fading protection."
    set_font(p3.runs[0], size=16, color=c_dark)
    p3.space_before = Pt(15)
    
    p4 = tf.add_paragraph()
    p4.text = "Demo Key Visuals:\n" \
              "• Green boxes: Centroid frequency estimates.\n" \
              "• Blue/Cyan trailing lines: Historical trajectory maps.\n" \
              "• Red highlights: Fading/Doppler null segments protected by track age."
    set_font(p4.runs[0], size=16, color=c_dark)
    p4.space_before = Pt(15)

    notes19 = slide19.notes_slide.notes_text_frame
    notes19.text = (
        "Now I will show you a live interactive demonstration of our tracking model.\n\n"
        "Our codebase includes a viewer.html file. By running the evaluation orchestrator, the system streams the "
        "incoming audio frames, runs the NMF decomposition and RL peak routing, and outputs JPEG frame summaries "
        "to output/frames.\n\n"
        "By opening viewer.html, we can watch the tracking decisions unfold at 10 frames per second, observing "
        "how the green trackers clamp onto vessel tonals and stay locked despite signal dropouts."
    )

    # Slide 20: Thank You & Questions Slide
    slide20 = add_slide("Thank You & Questions")
    tb = slide20.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.33), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p1 = tf.paragraphs[0]
    p1.text = "Thank You!"
    set_font(p1.runs[0], size=44, bold=True, color=c_navy)
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Questions?"
    set_font(p2.runs[0], size=32, bold=True, color=c_blue)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "\nCode Reference: https://github.com/roy-michael/rl-vessel-detection"
    set_font(p3.runs[0], size=18, color=c_gray)
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(30)

    notes20 = slide20.notes_slide.notes_text_frame
    notes20.text = (
        "Thank you all for your attention. I am now open to any questions you may have about the "
        "MDP formulation, reward design, NMF pre-processing, or algorithm performance."
    )

    # Save presentation
    output_path = "report/presentation_v4.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")

if __name__ == "__main__":
    create_presentation()
